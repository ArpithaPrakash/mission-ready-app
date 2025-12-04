import argparse
import json
import re
from pathlib import Path

try:
    from pptx import Presentation  # type: ignore
except Exception:  # pragma: no cover
    Presentation = None  # type: ignore

HEADINGS = [
    "MISSION", "PURPOSE", "INTENT", "COMMANDER'S INTENT",
    "KEY TASKS", "TASKS", "EXECUTION", "CONCEPT OF OPERATION",
    "SUSTAINMENT", "CMD AND SIGNAL", "COMMAND AND SIGNAL",
    "COMMAND & SIGNAL", "END STATE", "COORDINATING INSTRUCTIONS"
]


def extract_text_from_pptx(file_path: Path) -> str | None:
    """Extract text from all slides in a PPTX file."""
    if Presentation is None:
        print("python-pptx is not installed; cannot parse PPTX files.")
        return None
    try:
        presentation = Presentation(str(file_path))
        text_snippets: list[str] = []
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    shape_text = shape.text.strip()
                    if shape_text:
                        text_snippets.append(shape_text)
        return "\n".join(text_snippets)
    except Exception as error:
        print(f"Skipping file {file_path}: {error}")
        return None


def parse_conop_sections(text: str) -> dict[str, str]:
    """Split CONOP text into labeled sections using common headings."""
    if not text:
        return {}

    cleaned_text = re.sub(r"\n+", "\n", text)
    cleaned_text = re.sub(r"\s{2,}", " ", cleaned_text)

    pattern = r"(?i)(" + "|".join([re.escape(h) for h in HEADINGS]) + r")[:\s]*"
    parts = re.split(pattern, cleaned_text)

    sections: dict[str, str] = {}
    for index in range(1, len(parts), 2):
        heading = parts[index].strip().upper()
        body = parts[index + 1].strip()
        if body:
            sections[heading] = body
    return sections


def slugify(value: str) -> str:
    lowered = value.lower()
    lowered = re.sub(r"[^a-z0-9]+", "-", lowered)
    return re.sub(r"-+", "-", lowered).strip("-") or "conop"


def build_outpath(ppt_path: Path, outdir: Path) -> Path:
    outdir.mkdir(parents=True, exist_ok=True)
    stem = ppt_path.stem
    filename = f"{slugify(stem)}-conop.json"
    return outdir / filename


def process_pptx(ppt_path: Path, outdir: Path) -> bool:
    text = extract_text_from_pptx(ppt_path)
    if text is None:
        return False

    sections = parse_conop_sections(text)
    payload = {
        "filename": ppt_path.name,
        "path": str(ppt_path.resolve()),
        "sections": sections,
    }

    outpath = build_outpath(ppt_path, outdir)
    outpath.write_text(json.dumps(payload, ensure_ascii=False, indent=2), encoding="utf-8")
    print(f"Wrote: {outpath}")
    return True


def batch_process(input_dir: Path, outdir: Path) -> None:
    pptx_files = list(input_dir.rglob("*.pptx"))
    if not pptx_files:
        print(f"No PPTX files found in {input_dir}")
        return

    success = 0
    for ppt_path in pptx_files:
        if process_pptx(ppt_path, outdir):
            success += 1
    print(f"Successfully processed {success} of {len(pptx_files)} files")


def main() -> None:
    parser = argparse.ArgumentParser(description="Parse CONOP PPTX slides into structured JSON")
    parser.add_argument("input", help="Path to a PPTX file or directory containing PPTX files")
    parser.add_argument("--outdir", default="PARSED_CONOPS", help="Directory for JSON output")
    parser.add_argument("--batch", action="store_true", help="Process all PPTX files in the input directory")
    args = parser.parse_args()

    input_path = Path(args.input)
    outdir = Path(args.outdir)

    if not input_path.exists():
        print(f"Input not found: {input_path}")
        raise SystemExit(1)

    if args.batch or input_path.is_dir():
        batch_process(input_path, outdir)
    else:
        if not process_pptx(input_path, outdir):
            raise SystemExit(2)


if __name__ == "__main__":
    main()
